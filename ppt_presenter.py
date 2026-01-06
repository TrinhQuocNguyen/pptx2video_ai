#!/usr/bin/env python
# -*- coding: utf-8 -*-


import os
import tempfile
import argparse
from subprocess import call

from pdf2image import convert_from_path
from pptx import Presentation
from gtts import gTTS
import time
import asyncio
import edge_tts

import subprocess


# Define the path to your Poppler bin directory
# Download from https://github.com/oschwartz10612/poppler-windows/releases/
POPPLER_PATH = r'C:\Users\nguye\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin' # Update this path

# Define the path to your ffmpeg executable
# Download ffmpeg from https://github.com/BtbN/FFmpeg-Builds/releases
FFMPEG_PATH = r"C:\ffmpeg\bin\ffmpeg.exe"

__author__ = ['Nguyen Quoc Trinh']


## Sometimes ffmpeg is avconv
# FFMPEG_NAME = 'ffmpeg'
FFMPEG_NAME = 'avconv'

# Create a helper function for the voice
async def generate_voice(text, output_path):
    """Uses Edge-TTS for high-quality, stable Vietnamese speech."""
    communicate = edge_tts.Communicate(text, "vi-VN-HoaiMyNeural")
    await communicate.save(output_path)

def ppt_presenter1(pptx_path, pdf_path, output_path):
    with tempfile.TemporaryDirectory() as temp_path:
        images_from_path = convert_from_path(pdf_path, poppler_path=POPPLER_PATH)
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            if slide.has_notes_slide:
                print("Processing slide {}".format(i + 1))
                notes = slide.notes_slide.notes_text_frame.text
                image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
                image.save(image_path)
                # audio_path = os.path.join(temp_path, 'frame_{}.mp3'.format(i))

                # tts = gTTS(text=notes, lang='vi', slow=False)
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        audio_path = os.path.join(temp_path, f'frame_{i}.mp3')
                         # Run the voice generation
                        asyncio.run(generate_voice(notes, audio_path))
                # tts.save(audio_path)
                # time.sleep(2)  # Give Google a 2-second breather between slides


                ffmpeg_call(image_path, audio_path, temp_path, i)

        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i)) \
                      for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)

def ppt_presenter(pptx_path, pdf_path, output_path):
    # 1. Convert PDF to Images
    images = convert_from_path(pdf_path, poppler_path=POPPLER_PATH)
    prs = Presentation(pptx_path)
    
    if len(images) != len(prs.slides):
        print(f"Warning: PDF has {len(images)} pages but PPT has {len(prs.slides)} slides.")

    with tempfile.TemporaryDirectory() as temp_path:
        ts_files = []

        for i, (slide, image) in enumerate(zip(prs.slides, images)):
            print(f"Processing Slide {i+1}...")
            
            # Get notes or use a tiny bit of silence if empty
            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
            if not notes:
                notes = " " # Edge-TTS needs at least a space

            # Define paths
            img_p = os.path.join(temp_path, f"frame_{i}.jpg")
            aud_p = os.path.join(temp_path, f"frame_{i}.mp3")
            ts_p = os.path.join(temp_path, f"frame_{i}.ts")

            # Save Image and Audio
            image.save(img_p)
            asyncio.run(generate_voice(notes, aud_p))

            # Run FFmpeg to create a video segment (.ts is better for concatenating)
            
            
            # Replace your current 'cmd = [...]' block with this:
            cmd = [
                'C:\\ffmpeg\\bin\\ffmpeg.exe', 
                '-y', 
                '-loop', '1', 
                '-i', img_p, 
                '-i', aud_p,
                # THIS LINE IS THE FIX: It ensures the image is divisible by 2
                '-vf', "scale=trunc(iw/2)*2:trunc(ih/2)*2", 
                '-c:v', 'libx264', 
                '-c:a', 'aac', 
                '-shortest', 
                '-pix_fmt', 'yuv420p', 
                '-f', 'mpegts', 
                ts_p
            ]
            
            # Use subprocess.run to ensure it finishes before moving on
            subprocess.run(cmd, check=True)
            ts_files.append(ts_p)

        # 2. Concat all segments
        print("Finalizing video...")
        concat_list = "concat:" + "|".join(ts_files)
        final_cmd = [
            FFMPEG_PATH, "-y", "-i", concat_list, 
            "-c", "copy", "-bsf:a", "aac_adtstoasc", output_path
        ]
        subprocess.run(final_cmd, check=True)

    print(f"Done! Video saved to: {output_path}")

# Run it
# ppt_presenter("my_slides.pptx", "my_slides.pdf", "final_video.mp4")

def ffmpeg_call(image_path, audio_path, temp_path, i):
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))
    call([FFMPEG_PATH, '-loop', '1', '-y', '-i', image_path, '-i', audio_path,
          '-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'libfdk_aac',
          '-b:a', '192k', '-pix_fmt', 'yuv420p', '-shortest', out_path_mp4])
    call([FFMPEG_PATH, '-y', '-i', out_path_mp4, '-c', 'copy',
          '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts])


def ffmpeg_concat(video_list_str, out_path):
    call([FFMPEG_PATH, '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str),
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])


def main():
    parser = argparse.ArgumentParser(description='PPT Presenter help.')
    parser.add_argument('--pptx', help='input pptx path')
    parser.add_argument('--pdf', help='input pdf path')
    parser.add_argument('-o', '--output', help='output path')
    args = parser.parse_args()
    ppt_presenter(args.pptx, args.pdf, args.output)


if __name__ == '__main__':
    main()